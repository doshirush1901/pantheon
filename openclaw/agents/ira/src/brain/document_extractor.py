#!/usr/bin/env python3
"""
SHARED DOCUMENT EXTRACTOR - Single Source of Truth
===================================================

Consolidates 3 duplicate PDF extraction implementations:
- knowledge_retriever.py (pdfplumber)
- adaptive_retrieval.py (PyMuPDF/fitz)
- document_ingestor.py (pypdf)

Now all use this module with automatic fallback:
PyMuPDF → pdfplumber → pypdf

Usage:
    from document_extractor import DocumentExtractor, extract_pdf, extract_document
    
    # Simple API
    text = extract_pdf("/path/to/file.pdf")
    text = extract_document("/path/to/file.xlsx")
    
    # With caching
    extractor = DocumentExtractor()
    text = extractor.extract("/path/to/file.pdf")
"""

import hashlib
import os
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any, Union
from functools import lru_cache

# =============================================================================
# DATA CLASSES
# =============================================================================

@dataclass
class ExtractionResult:
    """Result of document extraction."""
    text: str
    page_count: int
    source: str
    extractor_used: str
    error: Optional[str] = None
    metadata: Dict[str, Any] = None
    
    def __post_init__(self):
        if self.metadata is None:
            self.metadata = {}
    
    @property
    def success(self) -> bool:
        return bool(self.text.strip()) and self.error is None


# =============================================================================
# PDF EXTRACTORS (with fallback chain)
# =============================================================================

def _extract_pdf_pymupdf(path: Path, max_pages: int = 100) -> ExtractionResult:
    """Extract PDF using PyMuPDF (fitz) - fastest option."""
    try:
        import fitz
        
        doc = fitz.open(path)
        text_parts = []
        page_count = len(doc)
        
        for i, page in enumerate(doc):
            if i >= max_pages:
                break
            text = page.get_text()
            if text.strip():
                text_parts.append(f"[PAGE {i+1}]\n{text}")
        
        doc.close()
        
        return ExtractionResult(
            text="\n\n".join(text_parts),
            page_count=page_count,
            source=str(path),
            extractor_used="PyMuPDF",
        )
    except ImportError:
        return ExtractionResult(
            text="", page_count=0, source=str(path),
            extractor_used="PyMuPDF", error="PyMuPDF not installed"
        )
    except Exception as e:
        return ExtractionResult(
            text="", page_count=0, source=str(path),
            extractor_used="PyMuPDF", error=str(e)
        )


def _extract_pdf_pdfplumber(path: Path, max_pages: int = 100) -> ExtractionResult:
    """Extract PDF using pdfplumber - good for tables."""
    try:
        import pdfplumber
        
        text_parts = []
        page_count = 0
        
        with pdfplumber.open(str(path)) as pdf:
            page_count = len(pdf.pages)
            for i, page in enumerate(pdf.pages):
                if i >= max_pages:
                    break
                text = page.extract_text()
                if text and text.strip():
                    text_parts.append(f"[PAGE {i+1}]\n{text}")
        
        return ExtractionResult(
            text="\n\n".join(text_parts),
            page_count=page_count,
            source=str(path),
            extractor_used="pdfplumber",
        )
    except ImportError:
        return ExtractionResult(
            text="", page_count=0, source=str(path),
            extractor_used="pdfplumber", error="pdfplumber not installed"
        )
    except Exception as e:
        return ExtractionResult(
            text="", page_count=0, source=str(path),
            extractor_used="pdfplumber", error=str(e)
        )


def _extract_pdf_pypdf(path: Path, max_pages: int = 100) -> ExtractionResult:
    """Extract PDF using pypdf - most compatible."""
    try:
        from pypdf import PdfReader
        
        reader = PdfReader(str(path))
        text_parts = []
        page_count = len(reader.pages)
        
        for i, page in enumerate(reader.pages):
            if i >= max_pages:
                break
            text = page.extract_text() or ""
            if text.strip():
                text_parts.append(f"[PAGE {i+1}]\n{text}")
        
        return ExtractionResult(
            text="\n\n".join(text_parts),
            page_count=page_count,
            source=str(path),
            extractor_used="pypdf",
        )
    except ImportError:
        return ExtractionResult(
            text="", page_count=0, source=str(path),
            extractor_used="pypdf", error="pypdf not installed"
        )
    except Exception as e:
        return ExtractionResult(
            text="", page_count=0, source=str(path),
            extractor_used="pypdf", error=str(e)
        )


def extract_pdf(path: Union[str, Path], max_pages: int = 100) -> str:
    """
    Extract text from PDF using best available library.
    
    Fallback chain: PyMuPDF → pdfplumber → pypdf
    
    Args:
        path: Path to PDF file
        max_pages: Maximum pages to extract (default 100)
    
    Returns:
        Extracted text content
    """
    path = Path(path)
    
    # Try extractors in order of preference
    for extractor in [_extract_pdf_pymupdf, _extract_pdf_pdfplumber, _extract_pdf_pypdf]:
        result = extractor(path, max_pages)
        if result.success:
            return result.text
    
    return ""


# =============================================================================
# OTHER DOCUMENT EXTRACTORS
# =============================================================================

def _extract_excel(path: Path) -> ExtractionResult:
    """Extract data from Excel files."""
    try:
        import pandas as pd
        
        # Read all sheets
        sheets = pd.read_excel(str(path), sheet_name=None)
        text_parts = []
        
        for sheet_name, df in sheets.items():
            if df.empty:
                continue
            
            text_parts.append(f"[SHEET: {sheet_name}]")
            
            # Convert headers
            headers = " | ".join(str(h) for h in df.columns)
            text_parts.append(f"Headers: {headers}")
            
            # Convert rows (limit to first 500 rows)
            for _, row in df.head(500).iterrows():
                row_text = " | ".join(str(v) for v in row.values if pd.notna(v))
                if row_text.strip():
                    text_parts.append(row_text)
            
            text_parts.append("")
        
        return ExtractionResult(
            text="\n".join(text_parts),
            page_count=len(sheets),
            source=str(path),
            extractor_used="pandas",
            metadata={"sheets": list(sheets.keys())}
        )
    except ImportError:
        return ExtractionResult(
            text="", page_count=0, source=str(path),
            extractor_used="pandas", error="pandas not installed"
        )
    except Exception as e:
        return ExtractionResult(
            text="", page_count=0, source=str(path),
            extractor_used="pandas", error=str(e)
        )


def _extract_docx(path: Path) -> ExtractionResult:
    """Extract text from Word documents."""
    try:
        import docx
        
        doc = docx.Document(str(path))
        text_parts = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)
        
        return ExtractionResult(
            text="\n\n".join(text_parts),
            page_count=len(doc.paragraphs) // 30 + 1,  # Approximate
            source=str(path),
            extractor_used="python-docx",
        )
    except ImportError:
        return ExtractionResult(
            text="", page_count=0, source=str(path),
            extractor_used="python-docx", error="python-docx not installed"
        )
    except Exception as e:
        return ExtractionResult(
            text="", page_count=0, source=str(path),
            extractor_used="python-docx", error=str(e)
        )


def _extract_text(path: Path) -> ExtractionResult:
    """Extract text from plain text files."""
    try:
        content = path.read_text(encoding='utf-8', errors='ignore')
        return ExtractionResult(
            text=content,
            page_count=1,
            source=str(path),
            extractor_used="text",
        )
    except Exception as e:
        return ExtractionResult(
            text="", page_count=0, source=str(path),
            extractor_used="text", error=str(e)
        )


def _extract_csv(path: Path) -> ExtractionResult:
    """Extract data from CSV files."""
    try:
        import pandas as pd
        
        df = pd.read_csv(str(path), nrows=1000)
        
        text_parts = []
        headers = " | ".join(str(h) for h in df.columns)
        text_parts.append(f"Columns: {headers}")
        text_parts.append("")
        
        for _, row in df.iterrows():
            row_text = " | ".join(str(v) for v in row.values if pd.notna(v))
            if row_text.strip():
                text_parts.append(row_text)
        
        return ExtractionResult(
            text="\n".join(text_parts),
            page_count=1,
            source=str(path),
            extractor_used="pandas-csv",
            metadata={"rows": len(df), "columns": list(df.columns)}
        )
    except ImportError:
        # Fallback to plain text reading
        return _extract_text(path)
    except Exception as e:
        return ExtractionResult(
            text="", page_count=0, source=str(path),
            extractor_used="pandas-csv", error=str(e)
        )


def _extract_pptx(path: Path) -> ExtractionResult:
    """Extract text from PowerPoint presentations."""
    try:
        from pptx import Presentation
        
        prs = Presentation(str(path))
        text_parts = []
        slide_count = len(prs.slides)
        
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
            
            if slide_texts:
                text_parts.append(f"[SLIDE {i+1}]")
                text_parts.extend(slide_texts)
                text_parts.append("")
        
        return ExtractionResult(
            text="\n".join(text_parts),
            page_count=slide_count,
            source=str(path),
            extractor_used="python-pptx",
            metadata={"slides": slide_count}
        )
    except ImportError:
        return ExtractionResult(
            text="", page_count=0, source=str(path),
            extractor_used="python-pptx", error="python-pptx not installed"
        )
    except Exception as e:
        return ExtractionResult(
            text="", page_count=0, source=str(path),
            extractor_used="python-pptx", error=str(e)
        )


# =============================================================================
# UNIFIED EXTRACTOR CLASS
# =============================================================================

class DocumentExtractor:
    """
    Unified document extraction with caching.
    
    Usage:
        extractor = DocumentExtractor()
        
        # Extract any supported document type
        result = extractor.extract("/path/to/document.pdf")
        print(result.text)
        
        # Check extraction success
        if result.success:
            process(result.text)
        else:
            print(f"Error: {result.error}")
    """
    
    EXTRACTORS = {
        ".pdf": lambda p: _extract_pdf_fallback(p),
        ".xlsx": _extract_excel,
        ".xls": _extract_excel,
        ".docx": _extract_docx,
        ".doc": _extract_docx,
        ".pptx": _extract_pptx,
        ".ppt": _extract_pptx,
        ".txt": _extract_text,
        ".md": _extract_text,
        ".csv": _extract_csv,
        ".json": _extract_text,
    }
    
    def __init__(self, cache_size: int = 100):
        """Initialize extractor with optional cache."""
        self._cache: Dict[str, ExtractionResult] = {}
        self._cache_size = cache_size
    
    def _cache_key(self, path: Path) -> str:
        """Generate cache key based on file path and modification time."""
        stat = path.stat()
        return f"{path}:{stat.st_mtime}:{stat.st_size}"
    
    def extract(self, path: Union[str, Path], use_cache: bool = True) -> ExtractionResult:
        """
        Extract text from a document.
        
        Args:
            path: Path to document
            use_cache: Whether to use cached results (default True)
        
        Returns:
            ExtractionResult with text and metadata
        """
        path = Path(path)
        
        if not path.exists():
            return ExtractionResult(
                text="", page_count=0, source=str(path),
                extractor_used="none", error=f"File not found: {path}"
            )
        
        # Check cache
        if use_cache:
            cache_key = self._cache_key(path)
            if cache_key in self._cache:
                return self._cache[cache_key]
        
        # Get extractor
        suffix = path.suffix.lower()
        extractor = self.EXTRACTORS.get(suffix)
        
        if extractor is None:
            return ExtractionResult(
                text="", page_count=0, source=str(path),
                extractor_used="none", error=f"Unsupported file type: {suffix}"
            )
        
        # Extract
        result = extractor(path)
        
        # Cache successful results
        if use_cache and result.success:
            cache_key = self._cache_key(path)
            self._cache[cache_key] = result
            
            # Trim cache if too large
            if len(self._cache) > self._cache_size:
                oldest = list(self._cache.keys())[0]
                del self._cache[oldest]
        
        return result
    
    def clear_cache(self):
        """Clear the extraction cache."""
        self._cache.clear()
    
    @property
    def supported_extensions(self) -> List[str]:
        """List of supported file extensions."""
        return list(self.EXTRACTORS.keys())


def _extract_pdf_fallback(path: Path) -> ExtractionResult:
    """Try all PDF extractors and return first successful result."""
    for extractor in [_extract_pdf_pymupdf, _extract_pdf_pdfplumber, _extract_pdf_pypdf]:
        result = extractor(path)
        if result.success:
            return result
    
    # Return last error if all failed
    return result


# =============================================================================
# CONVENIENCE FUNCTIONS
# =============================================================================

# Global extractor instance
_global_extractor: Optional[DocumentExtractor] = None


def get_extractor() -> DocumentExtractor:
    """Get or create global extractor instance."""
    global _global_extractor
    if _global_extractor is None:
        _global_extractor = DocumentExtractor()
    return _global_extractor


def extract_document(path: Union[str, Path]) -> str:
    """
    Extract text from any supported document type.
    
    Args:
        path: Path to document
    
    Returns:
        Extracted text content
    """
    result = get_extractor().extract(path)
    return result.text if result.success else ""


# =============================================================================
# EXPORTS
# =============================================================================

__all__ = [
    "ExtractionResult",
    "DocumentExtractor",
    "extract_pdf",
    "extract_document",
    "get_extractor",
]
